#!/usr/bin/python

import os
from pptx import Presentation
from pptx.util import Cm

frames_dir = "frames"

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

for filename in os.listdir(frames_dir):
    path = f'{frames_dir}/{filename}'
    print(path)
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = Cm(0)
    pic = slide.shapes.add_picture(path, left, top)

prs.save('bad-apple.pptx')